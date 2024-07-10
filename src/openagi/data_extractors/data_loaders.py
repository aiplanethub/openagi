import json
import tempfile
from tempfile import NamedTemporaryFile
from typing import Any
from urllib.parse import urlparse

import requests
from bs4 import BeautifulSoup as Soup
from langchain.docstore.document import Document
from langchain_community.document_loaders import (
    AirbyteStripeLoader,
    GitLoader,
    PyPDFLoader,
    RecursiveUrlLoader,
    TextLoader,
    UnstructuredMarkdownLoader,
    UnstructuredWordDocumentLoader,
    WebBaseLoader,
    YoutubeLoader,
)
from openagi.data_extractors.data_types import DataType
from openagi.data_extractors.data_source import DataSource

class DataLoader:
    def __init__(self, data_source: DataSource):
        self.data_source = data_source

    def load(self) -> Any:
        loader_methods = {
            DataType.TXT: self.load_txt,
            DataType.PDF: self.load_pdf,
            DataType.PPTX: self.load_pptx,
            DataType.DOCX: self.load_docx,
            DataType.GOOGLE_DOC: self.load_google_doc,
            DataType.MARKDOWN: self.load_markdown,
            DataType.GITHUB_REPOSITORY: self.load_github,
            DataType.WEBPAGE: self.load_webpage,
            DataType.YOUTUBE: self.load_youtube,
            DataType.URL: self.load_url,
        }

        loader = loader_methods.get(self.data_source.type)
        if loader:
            return loader()
        else:
            raise ValueError(f"Loader not implemented for type: {self.data_source.type}")

    def load_txt(self):
        with NamedTemporaryFile(suffix=".txt", delete=True) as temp_file:
            if self.data_source.url:
                file_response = requests.get(self.data_source.url).text
            else:
                file_response = self.data_source.content
            temp_file.write(file_response.encode())
            temp_file.flush()
            loader = TextLoader(file_path=temp_file.name)
            return loader.load_and_split()

    def load_pdf(self):
        if self.data_source.url:
            loader = PyPDFLoader(file_path=self.data_source.url)
        else:
            with NamedTemporaryFile(suffix=".pdf", delete=True) as temp_file:
                temp_file.write(self.data_source.content)
                temp_file.flush()
                loader = UnstructuredWordDocumentLoader(file_path=temp_file.name)
                return loader.load_and_split()
        return loader.load_and_split()

    def load_google_doc(self):
        pass

    def load_pptx(self):
        from pptx import Presentation

        with NamedTemporaryFile(suffix=".pptx", delete=True) as temp_file:
            if self.data_source.url:
                file_response = requests.get(self.data_source.url).content
            else:
                file_response = self.data_source.content
            temp_file.write(file_response)
            temp_file.flush()
            presentation = Presentation(temp_file.name)
            result = ""
            for i, slide in enumerate(presentation.slides):
                result += f"\n\nSlide #{i}: \n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        result += f"{shape.text}\n"
            return [Document(page_content=result)]

    def load_docx(self):
        with NamedTemporaryFile(suffix=".docx", delete=True) as temp_file:
            if self.data_source.url:
                file_response = requests.get(self.data_source.url).content
            else:
                file_response = self.data_source.content
            temp_file.write(file_response)
            temp_file.flush()
            loader = UnstructuredWordDocumentLoader(file_path=temp_file.name)
            return loader.load_and_split()

    def load_markdown(self):
        with NamedTemporaryFile(suffix=".md", delete=True) as temp_file:
            if self.data_source.url:
                file_response = requests.get(self.data_source.url).text
            else:
                file_response = self.data_source.content
            temp_file.write(file_response.encode())
            temp_file.flush()
            loader = UnstructuredMarkdownLoader(file_path=temp_file.name)
            return loader.load()

    def load_github(self):
        parsed_url = urlparse(self.data_source.url)
        path_parts = parsed_url.path.split("/")
        repo_name = path_parts[2]
        metadata = json.loads(self.data_source.metadata)

        with tempfile.TemporaryDirectory() as temp_dir:
            repo_path = f"{temp_dir}/{repo_name}/"
            loader = GitLoader(
                clone_url=self.data_source.url,
                repo_path=repo_path,
                branch=metadata["branch"],
            )
            return loader.load_and_split()

    def load_webpage(self):
        loader = RecursiveUrlLoader(
            url=self.data_source.url,
            max_depth=2,
            extractor=lambda x: Soup(x, "html.parser").text,
        )
        chunks = loader.load_and_split()
        for chunk in chunks:
            if "language" in chunk.metadata:
                del chunk.metadata["language"]
        return chunks

    def load_youtube(self):
        video_id = self.data_source.url.split("youtube.com/watch?v=")[-1]
        loader = YoutubeLoader(video_id=video_id)
        return loader.load_and_split()

    def load_url(self):
        url_list = self.data_source.url.split(",")
        loader = WebBaseLoader(url_list)
        return loader.load_and_split()